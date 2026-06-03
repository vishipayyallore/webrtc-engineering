"""PPTX → Markdown converter (raw extraction).

This script is meant to help *extract* slide text/tables/images into Markdown so that
Swamy can then *synthesize* proper learning notes in `src/weekN/01-notes/`,
`src/weekN/03-notebooks/`, and related week companions without violating the repo's
Zero-Copy Policy.

Design goals:
- No hardcoded paths (uses pathlib + CLI args)
- Safe defaults (writes outside source-material/ by default)
- Practical output for study workflows

Usage (PowerShell, from repo root):
    uv run python tools/pyscripts/pptx_to_md.py --input "C:\\path\\deck.pptx" --extract-images --include-notes

If you do not use uv:
    python tools/pyscripts/pptx_to_md.py --input "C:\\path\\deck.pptx"
"""

from __future__ import annotations

import argparse
import datetime as dt
import os
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape


def _is_within(child: Path, parent: Path) -> bool:
    """Return True if child is within parent (best-effort for Windows paths)."""
    try:
        child.resolve().relative_to(parent.resolve())
        return True
    except Exception:
        return False


def _safe_relpath(target: Path, start: Path) -> str:
    """Return a Markdown-friendly relative path (forward slashes)."""
    rel = os.path.relpath(target, start=start)
    return rel.replace("\\", "/")


def _extract_title(slide) -> str:
    # Prefer slide title placeholder if present
    if getattr(slide.shapes, "title", None) is not None and slide.shapes.title is not None:
        text = (slide.shapes.title.text or "").strip()
        if text:
            return text

    # Fallback: first non-empty text frame
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = (shape.text_frame.text or "").strip()
            if text:
                return text.splitlines()[0].strip()

    return "(Untitled)"


def _shape_text_to_markdown(shape: BaseShape) -> list[str]:
    lines: list[str] = []

    if not getattr(shape, "has_text_frame", False):
        return lines

    tf = shape.text_frame
    for paragraph in tf.paragraphs:
        raw = (paragraph.text or "").strip()
        if not raw:
            continue
        indent = "  " * int(getattr(paragraph, "level", 0) or 0)
        # Use bullet list formatting for simplicity.
        lines.append(f"{indent}- {raw}")

    return lines


def _table_to_markdown(shape: BaseShape, max_columns: int | None = None) -> list[str]:
    if not getattr(shape, "has_table", False):
        return []

    table = shape.table
    rows = table.rows
    cols = table.columns

    col_count = len(cols)
    if max_columns is not None:
        col_count = min(col_count, max_columns)

    def cell_text(r: int, c: int) -> str:
        text = (table.cell(r, c).text or "").strip()
        # Prevent Markdown table breakage
        return text.replace("|", "\\|")

    # Header: first row as header (common in slides)
    header = [cell_text(0, c) for c in range(col_count)]
    sep = ["---" for _ in range(col_count)]

    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(sep) + " |",
    ]

    for r in range(1, len(rows)):
        row = [cell_text(r, c) for c in range(col_count)]
        lines.append("| " + " | ".join(row) + " |")

    return lines


def _extract_images_from_slide(
    slide,
    images_dir: Path,
    slide_index: int,
    output_dir: Path,
) -> list[str]:
    lines: list[str] = []
    image_counter = 0

    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue

        image_counter += 1
        image = shape.image
        ext = (image.ext or "bin").lower()
        file_name = f"slide-{slide_index:02d}-image-{image_counter:02d}.{ext}"
        out_path = images_dir / file_name
        out_path.write_bytes(image.blob)

        rel = _safe_relpath(out_path, start=output_dir)
        lines.append(f"![]({rel})")

    return lines


def convert_pptx_to_markdown(
    input_pptx: Path,
    output_md: Path,
    extract_images: bool,
    images_dir: Path | None,
    include_notes: bool,
    omit_empty_slides: bool,
    max_table_columns: int | None,
) -> None:
    presentation = Presentation(str(input_pptx))

    output_dir = output_md.parent
    output_dir.mkdir(parents=True, exist_ok=True)

    if extract_images:
        if images_dir is None:
            images_dir = output_dir / f"{input_pptx.stem}_images"
        images_dir.mkdir(parents=True, exist_ok=True)

    now = dt.datetime.now().astimezone()

    lines: list[str] = []
    lines.append(f"# Raw PPTX Extract: {input_pptx.stem}")
    lines.append("")
    lines.append(f"- Source file: `{input_pptx.name}`")
    lines.append(f"- Extracted at: {now:%Y-%m-%d %H:%M %Z}")
    lines.append("")
    lines.append(
        "> This is an automated extraction from a PPTX for personal study. "
        "Please **synthesize** into original notes before publishing to `src/weekN/01-notes/` "
        "to respect the repository's Zero-Copy Policy."
    )

    for idx, slide in enumerate(presentation.slides, start=1):
        title = _extract_title(slide)
        slide_lines: list[str] = []

        slide_lines.append("")
        slide_lines.append(f"## Slide {idx}: {title}")
        slide_lines.append("")

        # Extract content: text frames + tables
        for shape in slide.shapes:
            table_md = _table_to_markdown(shape, max_columns=max_table_columns)
            if table_md:
                slide_lines.append("**Table**")
                slide_lines.extend(table_md)
                slide_lines.append("")

            text_md = _shape_text_to_markdown(shape)
            if text_md:
                slide_lines.extend(text_md)
                slide_lines.append("")

        # Images
        if extract_images and images_dir is not None:
            images_md = _extract_images_from_slide(
                slide=slide,
                images_dir=images_dir,
                slide_index=idx,
                output_dir=output_dir,
            )
            if images_md:
                slide_lines.append("**Images**")
                slide_lines.extend(images_md)
                slide_lines.append("")

        # Notes
        if include_notes and getattr(slide, "has_notes_slide", False):
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                slide_lines.append("<details>")
                slide_lines.append("<summary>Speaker notes</summary>")
                slide_lines.append("")
                for note_line in notes.splitlines():
                    slide_lines.append(note_line)
                slide_lines.append("")
                slide_lines.append("</details>")
                slide_lines.append("")

        # Optionally omit slides that contain only the header
        content_only = [ln for ln in slide_lines if ln.strip() and not ln.startswith("## Slide")]
        if omit_empty_slides and not content_only:
            continue

        lines.extend(slide_lines)

    output_md.write_text("\n".join(lines).rstrip() + "\n", encoding="utf-8")


def build_arg_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Convert PPTX slides to Markdown (raw extraction).")

    parser.add_argument(
        "--input",
        required=True,
        type=Path,
        help="Path to the .pptx file",
    )

    parser.add_argument(
        "--output",
        type=Path,
        default=None,
        help=(
            "Output Markdown path. Default: docs/exports/<pptx-stem>.md"
        ),
    )

    parser.add_argument(
        "--extract-images",
        action="store_true",
        help="Extract embedded images and link them in the Markdown.",
    )

    parser.add_argument(
        "--images-dir",
        type=Path,
        default=None,
        help="Directory to write extracted images (default: alongside output md).",
    )

    parser.add_argument(
        "--include-notes",
        action="store_true",
        help="Include speaker notes in collapsible <details> blocks.",
    )

    parser.add_argument(
        "--omit-empty-slides",
        action="store_true",
        help="Skip slides that produce no extracted content.",
    )

    parser.add_argument(
        "--max-table-columns",
        type=int,
        default=None,
        help="Limit the number of columns exported for wide tables.",
    )

    parser.add_argument(
        "--allow-source-material-output",
        action="store_true",
        help=(
            "Allow writing output under source-material/. "
            "(Not recommended; source-material is intended as read-only.)"
        ),
    )

    return parser


def main() -> int:
    parser = build_arg_parser()
    args = parser.parse_args()

    input_pptx: Path = args.input
    if input_pptx.suffix.lower() != ".pptx":
        raise SystemExit("Input must be a .pptx file")
    if not input_pptx.exists():
        raise SystemExit(f"Input file not found: {input_pptx}")

    workspace_root = Path(__file__).resolve().parents[1]
    source_material_dir = workspace_root / "source-material"

    output_md: Path
    if args.output is None:
        output_md = workspace_root / "docs" / "exports" / f"{input_pptx.stem}.md"
    else:
        output_md = args.output

    # Guardrails to respect repo policy.
    if _is_within(output_md, source_material_dir) and not args.allow_source_material_output:
        raise SystemExit(
            "Refusing to write output under source-material/. "
            "Pass --allow-source-material-output if you really want this."
        )

    images_dir = args.images_dir
    if images_dir is not None:
        if _is_within(images_dir, source_material_dir) and not args.allow_source_material_output:
            raise SystemExit(
                "Refusing to write images under source-material/. "
                "Pass --allow-source-material-output if you really want this."
            )

    convert_pptx_to_markdown(
        input_pptx=input_pptx,
        output_md=output_md,
        extract_images=bool(args.extract_images),
        images_dir=images_dir,
        include_notes=bool(args.include_notes),
        omit_empty_slides=bool(args.omit_empty_slides),
        max_table_columns=args.max_table_columns,
    )

    print(f"Wrote Markdown: {output_md}")
    if args.extract_images:
        actual_images_dir = images_dir or output_md.parent / f"{input_pptx.stem}_images"
        print(f"Extracted images (if any) to: {actual_images_dir}")

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
